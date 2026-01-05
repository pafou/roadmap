import json
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from styles_config import THEME_STYLING, TASK_STYLING, COLOR_MAPPING, BAR_TYPE_COLORS, MILESTONE_STYLES, MILESTONE_TEXT_STYLE

# --- Chargement du JSON ---
with open("data/roadmap.json", "r", encoding="utf-8") as file:
    roadmap_data = json.load(file)

# --- Template PowerPoint ---
prs = Presentation("data/Roadmap_template.pptx")
slide = prs.slides[0]

# --- Mapping mois → colonne ---
MOIS_INDEX = {
    "Jan": 0, "Fév": 1, "Mar": 2, "Avr": 3, "Mai": 4,
    "Juin": 5, "Juil": 6, "Aoû": 7, "Sep": 8, "Oct": 9,
    "Nov": 10, "Déc": 11
}

# --- Mapping thème → ligne (calculé dynamiquement) ---
# Extract unique themes from the YAML and create dynamic mapping
unique_themes = [theme["name"] for theme in roadmap_data["themes"]]
THEME_INDEX = {theme: idx for idx, theme in enumerate(unique_themes)}

# Calculate theme heights based on number of lines
# Each theme height = 0.3cm (title) + number_of_lines × 0.6cm (0.3cm bar + 0.3cm spacing)
theme_heights = {}
cumulative_heights = {}
current_y = Cm(2)  # Starting Y position (1.2 cm higher than original 3 cm)

for theme_data in roadmap_data["themes"]:
    theme_name = theme_data["name"]
    items = theme_data["items"]

    # Count the number of lines (each line object represents one line)
    num_lines = len(items)

    theme_height = Cm(0.3) + Cm(0.6) * num_lines  # 0.3cm title + lines × 0.6cm

    theme_heights[theme_name] = theme_height
    cumulative_heights[theme_name] = current_y
    current_y += theme_height

# Debug: show calculated themes and their heights
print("Thèmes calculés:", list(THEME_INDEX.keys()))
for theme_name, height in theme_heights.items():
    # Count lines - each item in the new structure represents a line
    num_lines = len(roadmap_data['themes'][THEME_INDEX[theme_name]]['items'])
    # Convert Cm object to actual centimeter value for display
    height_cm = height / 360000  # Convert from EMUs to cm (360000 EMUs = 1 cm)
    print(f"Thème '{theme_name}': {height_cm:.1f} cm ({num_lines} lignes)")

# --- Mise à jour des labels des thèmes dans le template ---
# Update theme labels in the PowerPoint template using configuration
for i, theme_name in enumerate(unique_themes):
    # Add theme formatting: title only (no horizontal line)
    # Theme title: 0.3cm high, starting 0.1cm from left
    title_y = cumulative_heights[theme_name]  # Position for this theme's title (dynamic)
    title_shape = slide.shapes.add_shape(
        1,  # rectangle shape for the title
        Cm(0.1), title_y,  # Start 0.1cm from left border
        Cm(5), Cm(0.3)  # 5cm wide, 0.3cm high
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
    title_shape.line.color.rgb = RGBColor(255, 255, 255)  # White border (invisible)
    title_shape.line.width = Pt(0)  # No border width

    # Add top border as a separate line shape
    top_border_y = title_y  # Position at the top of the title rectangle
    top_border_shape = slide.shapes.add_shape(
        1,  # line shape for top border
        Cm(0.1), top_border_y,  # Start 0.1cm from left border
        Cm(5), Cm(0)  # 5cm wide, 0 height (pure line)
    )
    top_border_shape.line.color.rgb = RGBColor(0x32, 0x4B, 0x6B)  # #324B6B color
    top_border_shape.line.width = Pt(1)  # 1pt thickness
    top_border_shape.fill.background()  # No fill for line

    # Add text to the title shape
    tf = title_shape.text_frame
    tf.text = theme_name
    for paragraph in tf.paragraphs:
        paragraph.font.name = THEME_STYLING['font_name']
        paragraph.font.size = Pt(THEME_STYLING['font_size'])
        paragraph.font.bold = THEME_STYLING['font_bold']
        paragraph.font.color.rgb = RGBColor(*THEME_STYLING['font_color'])
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT  # Left alignment

    # Find and update theme label shapes (assuming they exist in the template)
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text:
            # Look for placeholder text like "Thème 1", "Thème 2", etc.
            current_text = shape.text_frame.text.strip()
            if current_text.startswith("Thème ") and current_text[6:].isdigit():
                theme_number = int(current_text[6:])
                if theme_number == i + 1:  # theme numbers start from 1
                    shape.text_frame.text = theme_name
                    # Apply styling from configuration
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.name = THEME_STYLING['font_name']
                        paragraph.font.size = Pt(THEME_STYLING['font_size'])
                        paragraph.font.bold = THEME_STYLING['font_bold']
                        paragraph.font.color.rgb = RGBColor(*THEME_STYLING['font_color'])

def coord_x(mois):
    return Cm(3.6) + Cm(1.6) * MOIS_INDEX[mois]

def coord_y(theme, ligne):
    # Theme title is 0.3cm high, items start 0.3cm below theme title
    # Each line is 0.3cm bar height + 0.3cm spacing = 0.6cm total
    return cumulative_heights[theme] + Cm(0.3) + Cm(0.6) * (ligne - 1)

def get_color_for_type(item_type, subtype=None):
    """Get color based on item type and subtype"""
    if item_type == "bar" and subtype and subtype in BAR_TYPE_COLORS:
        return BAR_TYPE_COLORS[subtype]['background']
    elif item_type == "bar":
        return COLOR_MAPPING['charge_sure']
    elif item_type == "milestone":
        return COLOR_MAPPING['jalon']
    elif item_type == "text":
        return COLOR_MAPPING['default']
    else:
        return COLOR_MAPPING['default']

# --- Génération des shapes ---
for theme_data in roadmap_data["themes"]:
    theme_name = theme_data["name"]
    line_items = theme_data["items"]  # Each item is now a line object

    # Process each line and its items
    for line_index, line_item in enumerate(line_items):
        line_items_data = line_item["line"]["items"]
        line_number = line_index + 1  # Line numbers start from 1

        # Separate bars and milestones within this line
        bars = []
        milestones = []

        for item in line_items_data:
            if item["type"] == "bar":
                bars.append(item)
            elif item["type"] == "milestone":
                milestones.append(item)
            elif item["type"] == "text":
                # Handle text items if needed
                pass

        # Sort bars by start month to ensure proper ordering
        bars.sort(key=lambda x: MOIS_INDEX[x["start"]])

        # Draw bars on this line
        for bar in bars:
            start_x = coord_x(bar["start"])
            end_x = coord_x(bar["end"])
            width = end_x - start_x + Cm(2)  # Add full column width
            y = coord_y(theme_name, line_number)

            shape = slide.shapes.add_shape(
                1,  # rectangle
                start_x, y,
                width, Cm(0.3)
            )

            # couleur selon type et subtype
            fill = shape.fill
            fill.solid()
            bar_subtype = bar.get("subtype", None)
            fill.fore_color.rgb = RGBColor(*get_color_for_type(bar["type"], bar_subtype))

            # pas de bordure pour les barres
            shape.line.color.rgb = RGBColor(255, 255, 255)  # White border (invisible)
            shape.line.width = Pt(0)  # No border width
            shape.line.fill.background()  # No fill for the line

            # texte
            tf = shape.text_frame
            tf.text = bar["label"]
            tf.paragraphs[0].font.size = Pt(TASK_STYLING['font_size'])

            # Set text color to match foreground style for bar type
            if bar_subtype and bar_subtype in BAR_TYPE_COLORS:
                foreground_color = BAR_TYPE_COLORS[bar_subtype]['foreground']
                tf.paragraphs[0].font.color.rgb = RGBColor(*foreground_color)
            else:
                # Default text color if no specific foreground is defined
                tf.paragraphs[0].font.color.rgb = RGBColor(*THEME_STYLING['font_color'])

        # Draw milestones on this line
        for milestone in milestones:
            x = coord_x(milestone["month"])
            y = coord_y(theme_name, line_number) - Cm(0.3)  # Position 0.3cm higher as requested

            # Create a rectangle shape for milestone (text in rectangle, 0.3cm high, no background/border)
            shape = slide.shapes.add_shape(
                1,  # rectangle
                x + Cm(1.6), y,  # End of the column
                Cm(3), Cm(0.3)  # 2cm wide, 0.3cm high as requested
            )

            # No background (transparent) and no border as requested
            shape.fill.background()  # Transparent background
            shape.line.color.rgb = RGBColor(255, 255, 255)  # White border (invisible)
            shape.line.width = Pt(0)  # No border width
            shape.line.fill.background()  # Ensure line has no fill

            # texte
            tf = shape.text_frame
            tf.text = milestone["label"]
            # Apply MILESTONE_TEXT_STYLE
            tf.paragraphs[0].font.name = MILESTONE_TEXT_STYLE['font_name']
            tf.paragraphs[0].font.size = Pt(MILESTONE_TEXT_STYLE['font_size'])
            tf.paragraphs[0].font.color.rgb = RGBColor(*MILESTONE_TEXT_STYLE['font_color'])
            # Ensure text has no outline or border
            tf.paragraphs[0].font.outline = False
            tf.paragraphs[0].font.shadow = False
            tf.paragraphs[0].font.underline = False

            # Add downward-pointing triangle below the text rectangle
            # Triangle: 0.3cm height, 0.15cm width, positioned 0.1cm below the text rectangle
            triangle_y = y + Cm(0.3) + Cm(-0.1)  # 0.1cm below the text rectangle
            triangle_x = x + Cm(1.6) - Cm(0.075)  # Centered below the text rectangle

            # Create isosceles triangle shape (using proper MSO_AUTO_SHAPE_TYPE constant)
            triangle_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
                triangle_x, triangle_y,
                Cm(0.15), Cm(0.3)  # 0.15cm width, 0.3cm height
            )

            # Rotate triangle 180 degrees to make it point downward
            triangle_shape.rotation = 180  # degrees

            # Set triangle color based on milestone style
            milestone_style = milestone.get("style", "default")  # Default to "default" if no style specified
            if milestone_style in MILESTONE_STYLES:
                style_config = MILESTONE_STYLES[milestone_style]
                triangle_shape.fill.solid()
                triangle_shape.fill.fore_color.rgb = RGBColor(*style_config['triangle_background'])
                triangle_shape.line.color.rgb = RGBColor(*style_config['triangle_border'])
                triangle_shape.line.width = Pt(0.5)  # Thin border
            else:
                # Fallback to original behavior if style not found
                triangle_shape.fill.solid()
                triangle_shape.fill.fore_color.rgb = RGBColor(*get_color_for_type(milestone["type"]))
                triangle_shape.line.color.rgb = RGBColor(*get_color_for_type(milestone["type"]))
                triangle_shape.line.width = Pt(0.5)  # Thin border

prs.save("data/Roadmap_generee.pptx")
print("Roadmap générée !")
